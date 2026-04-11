"""This module contains utility functions that are used for creating the slides."""

import datetime
from lxml import etree
from moviepy import VideoFileClip
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor

# from pptx.enum.text import MSO_ANCHOR # vertical alignment of text
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from isda_pptgen.duplicate import duplicate_slide

# defining the template slides in the order they appear in the template
TEMPLATE = (
    "start_slide",
    "media_with_caption_slide",
    "simple_title_slide",
    "title_with_logo_slide",
    "welcome_and_announcements_slide",
    "membership_transfer_slide",
    "sermon_title_slide",
    "song_title_slide",
    "chorus_slide",
    "scripture_slide",
    "thithes_and_offerings_slide_0",
    "thithes_and_offerings_slide_1",
    "thithes_and_offerings_slide_2",
    "end_slide",
)

# debug
# print("Using structure of template :")
# for i, slide in enumerate(TEMPLATE):
#     print(f"{i} : {slide}")

COLORS = {
    "white": RGBColor(255, 255, 255),
    "pink": RGBColor(242, 207, 248),
    "green": RGBColor(217, 242, 208),
    "blue": RGBColor(166, 202, 236),
    "yellow": RGBColor(242, 219, 132),
}

FONTS = {
    "extrabold": "Nunito ExtraBold",
    "semibold": "Nunito SemiBold"
}

def clear_media_folder():
    """Clears the media folder of png, srt and mp4 files."""
    import os
    import glob
    extensions = ["*.png", "*.srt", "*.mp4"]
    files = []
    for ext in extensions:
        files.extend(glob.glob(f"media/{ext}"))
    
    for f in files:
        if os.path.isfile(f):
            os.remove(f)
    print("Media folder cleared of png, srt and mp4 files.")

def insert_text(pragraph, text:str, size:int, color:str="white", font:str="extrabold", superscript:bool=False):
    """
    Creates a 'run' (text with formatting) with the following settings:
    - font size: 24pt (-1), 32pt (0), 44pt (1), 60pt (2)
    - font color: white or rgb(242, 207, 248) (colored)
    - font name: Nunito ExtraBold ("extrabold") or Nunito SemiBold ("semibold")
    - alignment: center
    - superscript: False by default
    Does not return anything, modifies the paragraph object directly.
    """
    def set_subscript(font):
        font._element.set("baseline", "-25000")

    def set_superscript(font):
        font._element.set("baseline", "30000")

    def set_strikethrough(font):
        font._element.set("strike","sngStrike")

    text = str(text)

    match size:
        case -1:
            font_size = Pt(24)
        case 0:
            font_size = Pt(32)
        case 1:
            font_size = Pt(44)
        case 2:
            font_size = Pt(60)
        case _: # default
            font_size = Pt(60)

    custom_run = pragraph.add_run()
    custom_run.text = text
    custom_run.font.size = font_size
    custom_run.font.color.rgb = COLORS.get(color, COLORS["white"])
    custom_run.font.name = FONTS.get(font, FONTS["extrabold"])
    custom_run.alignment = PP_ALIGN.CENTER

    if superscript:
        set_superscript(custom_run.font)


def make_media_autoplay(media):
    """Sets the media to autoplay in a really ugly way that I don't understand but it works."""
    def xpath(el, query):
        """Helper function to find elements in the XML tree."""
        nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
        return etree.ElementBase.xpath(el, query, namespaces=nsmap)

    el_id = xpath(media.element, ".//p:cNvPr")[0].attrib["id"]
    el_cnt = xpath(
        media.element.getparent().getparent().getparent(),
        f'.//p:timing//p:video//p:spTgt[@spid="{el_id}"]',
    )[0]
    cond = xpath(el_cnt.getparent().getparent(), ".//p:cond")[0]
    cond.set("delay", "0")

def delete_template_slides(presentation:Presentation):
    """Deletes the template slides from the presentation. (the first ~7 slides)"""
    for _ in range(len(TEMPLATE)):
        presentation.slides._sldIdLst.remove(presentation.slides._sldIdLst[0])

def insert_start_slide(presentation:Presentation, date: datetime.date = None):
    """Inserts the start slide."""
    slide = duplicate_slide(presentation, TEMPLATE.index("start_slide"))
    
    if date:
        day = date.day
        if 11 <= (day % 100) <= 13:
            suffix = "th"
        else:
            suffix = {1: "st", 2: "nd", 3: "rd"}.get(day % 10, "th")
        
        month_year = date.strftime(" %B %Y")
        
        for shape in slide.shapes:
            if shape.name == "date":
                shape.text_frame.clear()
                date_paragraph = shape.text_frame.paragraphs[0]
                date_paragraph.alignment = PP_ALIGN.RIGHT
                insert_text(date_paragraph, str(day), size=0, color="white")
                insert_text(date_paragraph, suffix, size=0, color="white", superscript=True)
                insert_text(date_paragraph, month_year, size=0, color="white")

def insert_video_slide(presentation:Presentation, video_path:str, thumbnail_path:str, caption:str=""):
    """Inserts a slide with a video and a caption."""
    slide = duplicate_slide(presentation, TEMPLATE.index("media_with_caption_slide"))
    caption = caption.strip()

    # getting size and aspect ratio of the video
    clip = VideoFileClip(video_path)
    aspect_ratio = clip.size[0] / clip.size[1]
    height = presentation.slide_height.inches
    width = presentation.slide_height.inches * aspect_ratio
    left = (presentation.slide_width.inches - width) / 2

    # changing the caption
    if caption != "":
        for shape in slide.shapes:
            if shape.name == "caption":
                shape.text_frame.clear()
                caption_paragraph = shape.text_frame.paragraphs[0]
                caption_paragraph.alignment = PP_ALIGN.CENTER
                insert_text(caption_paragraph, caption, size=-1, color="white")
    else:
        # removing the caption shape if there's no caption
        for shape in slide.shapes:
            if shape.name == "caption":
                slide.shapes._spTree.remove(shape._element)

    # addding the video
    video = slide.shapes.add_movie(video_path, Inches(left), 0, Inches(width), Inches(height), poster_frame_image=thumbnail_path)
    slide.shapes._spTree.remove(video._element)
    slide.shapes._spTree.insert(2, video._element)
    make_media_autoplay(video)

def insert_image_slide(presentation:Presentation, image_path:str, caption:str=""):
    """Inserts a slide with an image and a caption."""
    slide = duplicate_slide(presentation, TEMPLATE.index("media_with_caption_slide"))
    caption = caption.strip()

    # getting size and aspect ratio of the image
    img = Image.open(image_path)
    width, height = img.size
    aspect_ratio = width / height
    height = presentation.slide_height.inches
    width = presentation.slide_height.inches * aspect_ratio
    left = (presentation.slide_width.inches - width) / 2

    # changing the caption
    if caption != "":
        for shape in slide.shapes:
            if shape.name == "caption":
                shape.text_frame.clear()
                caption_paragraph = shape.text_frame.paragraphs[0]
                caption_paragraph.alignment = PP_ALIGN.CENTER
                insert_text(caption_paragraph, caption, size=0, color="white")
    else:
        # removing the caption shape if there's no caption
        for shape in slide.shapes:
            if shape.name == "caption":
                slide.shapes._spTree.remove(shape._element)

    # adding the image
    image = slide.shapes.add_picture(image_path, Inches(left), 0, Inches(width), Inches(height))
    slide.shapes._spTree.remove(image._element)
    slide.shapes._spTree.insert(2, image._element)

def insert_images(presentation:Presentation, images:tuple, caption:str=""):
    """Batch inserts images into the presentation. Each image will be on a separate slide. The caption will be the same for all the images."""
    for image in images:
        insert_image_slide(presentation, image, caption)

def insert_simple_title_slide(presentation:Presentation, title:str):
    """Inserts a slide with a simple title."""
    slide = duplicate_slide(presentation, TEMPLATE.index("simple_title_slide"))

    # changing the contents of the slide
    for shape in slide.shapes:
        if shape.name == "title":
            shape.text_frame.clear()
            title_paragraph = shape.text_frame.paragraphs[0]
            title_paragraph.alignment = PP_ALIGN.CENTER
            insert_text(title_paragraph, title, size=2, color="white")

def insert_title_with_logo_slide(presentation:Presentation, title:str):
    """Inserts a slide with a title and a logo."""
    slide = duplicate_slide(presentation, TEMPLATE.index("title_with_logo_slide"))

    # changing the contents of the slide
    for shape in slide.shapes:
        if shape.name == "title":
            shape.text_frame.clear()
            title_paragraph = shape.text_frame.paragraphs[0]
            title_paragraph.alignment = PP_ALIGN.LEFT
            insert_text(title_paragraph, title, size=2, color="white")

def insert_welcome_and_announcements_slide(presentation:Presentation):
    """Inserts the welcome and announcements slide."""
    duplicate_slide(presentation, TEMPLATE.index("welcome_and_announcements_slide"))

def insert_membership_transfer_slide(presentation:Presentation, in_or_out:str, reading:str, name:str, church:str):
    """Inserts a slide for a membership transfer.
    - in_or_out: "in" or "out"
    - reading: the reading for the transfer (e.g. "first", "second",
    - name: the name of the person transferring
    - church: the church they are transferring from/to
    """
    in_or_out = in_or_out.lower()
    if in_or_out not in ["in", "out"]:
        raise ValueError("in_or_out must be either 'in' or 'out'")
    from_or_to = "from" if in_or_out == "in" else "to"
    slide = duplicate_slide(presentation, TEMPLATE.index("membership_transfer_slide"))

    # changing the contents of the slide
    for shape in slide.shapes:
        if shape.name == "text":
            shape.text_frame.clear()
            text_paragraph = shape.text_frame.paragraphs[0]
            text_paragraph.alignment = PP_ALIGN.CENTER
            insert_text(text_paragraph, f"Membership Transfer {in_or_out}\n", size=2, color="white", font="semibold")
            insert_text(text_paragraph, f"{reading.capitalize()} reading\n\n", size=2, color="pink", font="semibold")
            insert_text(text_paragraph, f"{name}\n\n", size=2, color="yellow")
            insert_text(text_paragraph, f"{from_or_to} ", size=2, color="white", font="semibold")
            insert_text(text_paragraph, church, size=2, color="green", font="semibold")


def insert_sermon_title_slide(presentation:Presentation, sermon_title:str, preacher:str=""):
    """Inserts a slide with the title of the sermon."""
    slide = duplicate_slide(presentation, TEMPLATE.index("sermon_title_slide"))

    # changing the contents of the slide
    for shape in slide.shapes:
        if shape.name == "title":
            shape.text_frame.clear()
            title_paragraph = shape.text_frame.paragraphs[0]
            title_paragraph.alignment = PP_ALIGN.LEFT
            insert_text(title_paragraph, "Sermon : \n", size=2, color="blue")
            insert_text(title_paragraph, sermon_title, size=2, color="white")
        elif shape.name == "preacher":
            shape.text_frame.clear()
            preacher_paragraph = shape.text_frame.paragraphs[0]
            preacher_paragraph.alignment = PP_ALIGN.LEFT
            insert_text(preacher_paragraph, preacher, size=1, color="white")

def insert_song_title_slide(presentation:Presentation, number:int, title:str):
    """Inserts a slide with the title of the song."""
    number = str(f"{number:03}")
    slide = duplicate_slide(presentation, TEMPLATE.index("song_title_slide"))

    # changing the contents of the slide
    for shape in slide.shapes:
        if shape.name == "song_title":
            shape.text_frame.clear()
            title_paragraph = shape.text_frame.paragraphs[0]
            title_paragraph.alignment = PP_ALIGN.CENTER
            insert_text(title_paragraph, f"#{number}", size=0, color="white")
            insert_text(title_paragraph, f"\n{title}", size=2, color="white")

def insert_chorus_slide(presentation:Presentation, verse_name:str, text:str, last_slide:bool=False):
    """Inserts a slide with the chorus of the song."""
    slide = duplicate_slide(presentation, TEMPLATE.index("chorus_slide"))

    # changing the contents of the slide
    for shape in slide.shapes:
        if shape.name == "chorus":
            shape.text_frame.clear()
            chorus_paragraph = shape.text_frame.paragraphs[0]
            chorus_paragraph.alignment = PP_ALIGN.CENTER
            if verse_name != "":
                insert_text(chorus_paragraph, verse_name, size=0, color="pink")
                insert_text(chorus_paragraph, f"\n{text}", size=1, color="white")
            else:
                insert_text(chorus_paragraph, text, size=1, color="white")

    if last_slide:
        # get size of the slide in inches
        posx = Inches(presentation.slide_width.inches) - Inches(0.8)
        posy = Inches(presentation.slide_height.inches) - Inches(0.8)

        # add the square
        square = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, posx, posy, Inches(0.4), Inches(0.4))
        square.fill.solid()
        square.fill.fore_color.rgb = RGBColor(255, 255, 255)
        # remove outline
        square.line.fill.background()

def insert_smart_chorus_slide(presentation:Presentation, verse_name:str, text:str, last_slide:bool=False):
    """Same as insert_chorus_slide but will split the text into multiple slides if it's too long."""
    MAX_LINES = 8
    lines = text.split("\n")

    # attempt manual split (some verses have an empty line where it makes sence that they would be split)
    manual_split = text.split("\n\n")
    if len(manual_split) > 1:
        insert_chorus_slide(presentation, verse_name, manual_split[0]) # verse_name only on the first slide
        for i in range(1, len(manual_split)):
            insert_chorus_slide(presentation, "", manual_split[i], last_slide if i == len(manual_split)-1 else False)
        return


    if len(lines) <= MAX_LINES:
        insert_chorus_slide(presentation, verse_name, text, last_slide)
    else:
        # attempt automatic splitting : splitting into multiple slides by splitting the lines evenly

        lines_per_slide_to_be_even = len(lines)
        while lines_per_slide_to_be_even > MAX_LINES:
            lines_per_slide_to_be_even = int(lines_per_slide_to_be_even / 2)
            # this will divide the lines into 2 parts until the number of lines is less than or equal to MAX_LINES

        # verse_name only on the first slide
        insert_chorus_slide(presentation, verse_name, "\n".join(lines[:lines_per_slide_to_be_even]))

        for i in range(lines_per_slide_to_be_even, len(lines), lines_per_slide_to_be_even):
            # if it's the last slide set last_slide to whatever is provided
            if i + lines_per_slide_to_be_even >= len(lines):
                insert_chorus_slide(presentation, "", "\n".join(lines[i:i+lines_per_slide_to_be_even]), last_slide)
            # if the next slide will have less than (lines_per_slide_to_be_even) lines, add the lines from that slide to this one
            elif len(lines[i+lines_per_slide_to_be_even:i+lines_per_slide_to_be_even+lines_per_slide_to_be_even]) < lines_per_slide_to_be_even:
                insert_chorus_slide(presentation, "", "\n".join(lines[i:]), last_slide)
                break
            else:
                insert_chorus_slide(presentation, "", "\n".join(lines[i:i+lines_per_slide_to_be_even]))

def insert_scripture_slide(presentation:Presentation, reference:str, text:tuple, verse_separator:str=" "):
    """
    Inserts a slide with the verses.
    Text need to be provided like this:
    (
        ("1", "This is the first verse."),
        ("2", "This is the second verse."),
        ("3", "This is the third verse.")
    )
    """
    slide = duplicate_slide(presentation, TEMPLATE.index("scripture_slide"))

    # changing the contents of the slide
    for shape in slide.shapes:
        if shape.name == "scripture":
            shape.text_frame.clear()
            verse_paragraph = shape.text_frame.paragraphs[0]
            verse_paragraph.alignment = PP_ALIGN.CENTER
            for verse in text:
                if verse[0] != "":
                    if verse[0] != text[0][0]:
                        insert_text(verse_paragraph, verse_separator, size=1, color="white")
                    insert_text(verse_paragraph, verse[0]+" ", size=1, color="pink", superscript=True)
                insert_text(verse_paragraph, verse[1], size=1, color="white")
        elif shape.name == "reference":
            shape.text_frame.clear()
            reference_paragraph = shape.text_frame.paragraphs[0]
            reference_paragraph.alignment = PP_ALIGN.CENTER
            insert_text(reference_paragraph, reference, size=0, color="pink")

def insert_thithes_and_offerings_slides(presentation:Presentation, unallocated_offerings:str):
    """Inserts the slides for the thithes and offerings given the unallocated offerings."""
    offering_slide_0 = duplicate_slide(presentation, TEMPLATE.index("thithes_and_offerings_slide_0"))

    # changing the contents of the slide
    for shape in offering_slide_0.shapes:
        if shape.name == "unallocated_offerings":
            shape.text_frame.clear()
            reference_paragraph = shape.text_frame.paragraphs[0]
            reference_paragraph.alignment = PP_ALIGN.CENTER
            insert_text(reference_paragraph, "Today’s unallocated offerings will go towards ", size=0, color="white", font="semibold")
            insert_text(reference_paragraph, unallocated_offerings, size=0, color="green", font="semibold")
    
    offering_slide_1 = duplicate_slide(presentation, TEMPLATE.index("thithes_and_offerings_slide_1"))

    for shape in offering_slide_1.shapes:
        if shape.name == "unallocated_offerings":
            shape.text_frame.clear()
            reference_paragraph = shape.text_frame.paragraphs[0]
            reference_paragraph.alignment = PP_ALIGN.CENTER
            insert_text(reference_paragraph, "Today’s unallocated offerings will go towards ", size=-1, color="white", font="semibold")
            insert_text(reference_paragraph, unallocated_offerings, size=-1, color="green", font="semibold")
    
    offering_slide_2 = duplicate_slide(presentation, TEMPLATE.index("thithes_and_offerings_slide_2"))

    for shape in offering_slide_2.shapes:
        if shape.name == "unallocated_offerings":
            shape.text_frame.clear()
            reference_paragraph = shape.text_frame.paragraphs[0]
            reference_paragraph.alignment = PP_ALIGN.CENTER
            insert_text(reference_paragraph, "Today’s unallocated offerings will go towards ", size=-1, color="white", font="semibold")
            insert_text(reference_paragraph, unallocated_offerings, size=-1, color="green", font="semibold")

    insert_chorus_slide(presentation, "Doxology", "Praise God, from whom all blessings flow\nPraise Him, all creatures here below\nPraise Him above, ye heavenly host\nPraise Father, Son and Holy Ghost\nAmen!")

def insert_end_slide(presentation:Presentation):
    """Inserts the end slide."""
    duplicate_slide(presentation, TEMPLATE.index("end_slide"))

def insert_hymn(presentation:Presentation, number:int):
    """Inserts slides for a hymn from the JSON database."""
    import json
    
    with open("assets/hymns.json", "r", encoding="utf-8") as f:
        hymns_data = json.load(f)
        
    # Find the requested hymn
    target_hymn = None
    for hymn in hymns_data:
        if hymn["id"] == number:
            target_hymn = hymn
            break
            
    if not target_hymn:
        print(f"Hymn {number} not found!")
        return

    hymn_name = target_hymn["name"]
    lyrics = target_hymn.get("lyrics", [])

    insert_song_title_slide(presentation, number, hymn_name)
    
    num_blocks = len(lyrics)
    for i, block in enumerate(lyrics):
        is_last = (i == num_blocks - 1)
        b_type = block.get("type", "unknown")
        b_label = block.get("label", "")
        b_text = block.get("text", "")
        
        # Present title based on type and label
        if b_type == "verse":
            title = str(b_label)
        elif b_type == "refrain":
            title = "Refrain" if not b_label else b_label
        elif b_type == "bridge":
            title = "Bridge" if not b_label else b_label
        else:
            title = b_label

        insert_smart_chorus_slide(presentation, title, b_text, is_last)
