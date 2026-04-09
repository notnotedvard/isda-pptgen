"""Utilities to merge .pptx files into an existing presentation.

Strategy: preserve source formatting as faithfully as possible by cloning the
source slide masters and layouts into the target before copying slides.  This
mirrors what PowerPoint does internally when you "Keep Source Formatting" during
a manual copy-paste.
"""

from __future__ import annotations

import copy
import logging
from typing import Dict, Tuple

from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

from isda_pptgen.duplicate import (
    _object_rels,
    _exp_add_slide,
    clone_slide_layout,
    clone_slide_master,
    copy_shapes,
    remove_shape,
)

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------


def _master_key(master) -> str:
    """Return a stable fingerprint for a slide master based on its theme name."""
    try:
        names = master.element.xpath(
            ".//p:txStyles/p:titleStyle/a:lstStyle/../../../../../../"
            "following-sibling::p:cSld/@name",
            namespaces={
                "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
                "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            },
        )
    except Exception:
        names = []

    # Fall back to theme name inside the master
    try:
        theme_names = master.element.xpath(
            ".//a:theme/@name",
            namespaces={
                "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            },
        )
    except Exception:
        theme_names = []

    key = "|".join(names + theme_names)
    # Last resort: use XML hash
    if not key:
        key = str(hash(etree.tostring(master.element)))
    return key


def _layout_key(layout) -> str:
    """Return a stable fingerprint for a slide layout."""
    try:
        name = layout.name or ""
    except Exception:
        name = ""
    try:
        xml_hash = str(hash(etree.tostring(layout.element)))
    except Exception:
        xml_hash = ""
    return f"{name}|{xml_hash}"


def _copy_all_rels(src_part, dest_part) -> None:
    """Copy all relationships from src_part to dest_part (best-effort)."""
    for rel in _object_rels(src_part):
        try:
            if rel.is_external:
                dest_part.rels.get_or_add_ext_rel(rel.reltype, rel._target)
            else:
                dest_part.rels.get_or_add(rel.reltype, rel._target)
        except Exception as exc:
            logger.debug("Skipping rel %s: %s", rel.reltype, exc)


def _copy_slide_rels(src_slide, dest_slide) -> None:
    """Copy slide-level relationships (hyperlinks, audio, video, images, …)."""
    # Relationship types that are safe to clone at the slide level.
    # We intentionally skip slideLayout / slideMaster because those are
    # managed separately.
    skip_reltypes = {
        RT.SLIDE_LAYOUT,
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout",
    }
    for rel in _object_rels(src_slide.part):
        if rel.reltype in skip_reltypes:
            continue
        try:
            if rel.is_external:
                dest_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel._target)
            else:
                dest_slide.part.rels.get_or_add(rel.reltype, rel._target)
        except Exception as exc:
            logger.debug("Skipping slide rel %s: %s", rel.reltype, exc)


def _copy_slide_background(src_slide, dest_slide) -> None:
    """Copy an explicit slide background (p:bg element) if present."""
    src_el = src_slide.element
    dest_el = dest_slide.element

    # p:bg is a direct child of p:cSld
    src_cSld = src_el.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}cSld"
    )
    dest_cSld = dest_el.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}cSld"
    )

    if src_cSld is None or dest_cSld is None:
        return

    bg_tag = "{http://schemas.openxmlformats.org/presentationml/2006/main}bg"
    src_bg = src_cSld.find(bg_tag)
    if src_bg is None:
        return

    # Remove existing background on destination
    dest_bg = dest_cSld.find(bg_tag)
    if dest_bg is not None:
        dest_cSld.remove(dest_bg)

    dest_cSld.insert(0, copy.deepcopy(src_bg))


def _copy_slide_transition(src_slide, dest_slide) -> None:
    """Copy transition element (p:transition) if present."""
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    trans_tag = f"{{{p_ns}}}transition"

    src_trans = src_slide.element.find(trans_tag)
    if src_trans is None:
        return

    dest_trans = dest_slide.element.find(trans_tag)
    if dest_trans is not None:
        dest_slide.element.remove(dest_trans)

    dest_slide.element.append(copy.deepcopy(src_trans))


def _copy_slide_timing(src_slide, dest_slide) -> None:
    """Copy timing/animation element (p:timing) if present."""
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    timing_tag = f"{{{p_ns}}}timing"

    src_timing = src_slide.element.find(timing_tag)
    if src_timing is None:
        return

    dest_timing = dest_slide.element.find(timing_tag)
    if dest_timing is not None:
        dest_slide.element.remove(dest_timing)

    dest_slide.element.append(copy.deepcopy(src_timing))


# ---------------------------------------------------------------------------
# Master / layout cloning with deduplication
# ---------------------------------------------------------------------------


def _import_masters(
    target: Presentation,
    source: Presentation,
) -> Tuple[Dict[str, object], Dict[str, object]]:
    """
    Clone every slide master (and its layouts) from *source* into *target*.

    Returns two dicts:
      - master_map  : source master key  → cloned target master object
      - layout_map  : source layout key  → cloned target layout object
    """
    master_map: Dict[str, object] = {}
    layout_map: Dict[str, object] = {}

    for src_master in source.slide_masters:
        mkey = _master_key(src_master)

        # Check whether an equivalent master already exists in target
        existing_master = None
        for tm in target.slide_masters:
            if _master_key(tm) == mkey:
                existing_master = tm
                break

        if existing_master is not None:
            cloned_master = existing_master
            logger.debug("Reusing existing master: %s", mkey)
        else:
            cloned_master = clone_slide_master(target, src_master)
            logger.debug("Cloned new master: %s", mkey)

        master_map[mkey] = cloned_master

        # Clone layouts for this master
        for src_layout in src_master.slide_layouts:
            lkey = _layout_key(src_layout)

            # Reuse if already present on the cloned master
            existing_layout = None
            for tl in cloned_master.slide_layouts:
                if _layout_key(tl) == lkey:
                    existing_layout = tl
                    break

            if existing_layout is not None:
                layout_map[lkey] = existing_layout
                logger.debug("Reusing existing layout: %s", lkey)
            else:
                clone_slide_layout(target, src_layout, cloned_master)
                new_layout = cloned_master.slide_layouts[-1]
                layout_map[lkey] = new_layout
                logger.debug("Cloned new layout: %s", lkey)

    return master_map, layout_map


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def merge_pptx(target: Presentation, source_path: str) -> None:
    """
    Merge all slides from *source_path* into *target*, preserving the source
    formatting as faithfully as possible.

    The approach:
    1. Clone source slide masters & layouts into target (deduplicating by key).
    2. For each source slide, add a new slide backed by the *cloned* layout that
       matches the source slide's original layout.
    3. Copy all shapes, background, transitions, animations, notes, and
       slide-level relationships.

    :param target:      The target :class:`Presentation` to merge into.
    :param source_path: Path to the source .pptx file whose slides to append.
    """
    source = Presentation(source_path)

    # Step 1 — import all masters/layouts from source into target
    _master_map, layout_map = _import_masters(target, source)

    # Fallback layout (blank) in case a match cannot be found
    fallback_layout = target.slide_layouts[6]

    # Step 2 — copy each slide
    for slide_idx, src_slide in enumerate(source.slides):
        # Find the matching cloned layout
        src_layout = src_slide.slide_layout
        lkey = _layout_key(src_layout)
        dest_layout = layout_map.get(lkey, fallback_layout)

        # Add a new blank slide using the correct layout
        dest_slide = _exp_add_slide(target, dest_layout)

        # Remove placeholder shapes injected by the layout clone
        for shape in list(dest_slide.shapes):
            remove_shape(shape)

        # Copy shapes (handles groups, images, charts, tables, text, etc.)
        copy_shapes(src_slide.shapes, dest_slide)

        # Copy explicit slide background (overrides layout/master background)
        _copy_slide_background(src_slide, dest_slide)

        # Copy slide-level relationships (hyperlinks, media, …)
        _copy_slide_rels(src_slide, dest_slide)

        # Copy transitions and animations
        _copy_slide_transition(src_slide, dest_slide)
        _copy_slide_timing(src_slide, dest_slide)

        # Copy speaker notes
        if src_slide.has_notes_slide:
            try:
                src_notes_tf = src_slide.notes_slide.notes_text_frame
                dest_notes_tf = dest_slide.notes_slide.notes_text_frame
                # Deep-copy the XML tree of the text body for full formatting
                src_txBody = src_notes_tf._txBody
                dest_txBody = dest_notes_tf._txBody
                parent = dest_txBody.getparent()
                idx = list(parent).index(dest_txBody)
                parent.remove(dest_txBody)
                parent.insert(idx, copy.deepcopy(src_txBody))
            except Exception as exc:
                logger.warning("Could not copy notes for slide %d: %s", slide_idx, exc)
