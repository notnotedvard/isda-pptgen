"""This module contains utility functions that are used for creating the slides."""

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
# from pptx.enum.text import MSO_ANCHOR # vertical alignment of text
from pptx.enum.shapes import MSO_SHAPE
from moviepy.editor import VideoFileClip
from lxml import etree
from PIL import Image
from copy_slides import duplicate_slide

# defining the template slides in the order they appear in the template
TEMPLATE = (
    "logo_slide",
    "media_with_caption_slide",
    "simple_title_slide",
    "song_title_slide",
    "chorus_slide",
    "scripture_slide",
    "end_slide",
)

def insert_text(pragraph, text:str, size:int, colored:bool=False, superscript:bool=False):
    """Creates a 'run' (text with formatting) with the following settings:
    - font size: 32pt (0), 44pt (1), 60pt (2)
    - font color: white or rgb(242, 207, 248) (colored)
    - font name: Loos Normal Medium
    - alignment: center
    - superscript: False by default
    Does not return anything, modifies the paragraph object directly.
    """
    def set_subscript(font):
        font._element.set('baseline', '-25000')

    def set_superscript(font):
        font._element.set('baseline', '30000')

    def set_strikethrough(font):
        font._element.set('strike','sngStrike')

    text = str(text)

    match size:
        case 0:
            font_size = Pt(32)
        case 1:
            font_size = Pt(44)
        case 2:
            font_size = Pt(60)
        case _: # default
            font_size = Pt(60)

    custom_run = pragraph.add_run()
    custom_run.text = text
    custom_run.font.size = font_size
    custom_run.font.color.rgb = RGBColor(242, 207, 248) if colored else RGBColor(255, 255, 255)
    custom_run.font.name = "Loos Normal Medium"
    custom_run.alignment = PP_ALIGN.CENTER

    if superscript:
        set_superscript(custom_run.font)


def autoplay_media(media):
    """Sets the media to autoplay in a really ugly way that I don't understand but it works."""
    def xpath(el, query):
        """Helper function to find elements in the XML tree."""
        nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        return etree.ElementBase.xpath(el, query, namespaces=nsmap)
    
    el_id = xpath(media.element, './/p:cNvPr')[0].attrib['id']
    el_cnt = xpath(
        media.element.getparent().getparent().getparent(),
        './/p:timing//p:video//p:spTgt[@spid="%s"]' % el_id,
    )[0]
    cond = xpath(el_cnt.getparent().getparent(), './/p:cond')[0]
    cond.set('delay', '0')

def delete_template_slides(presentation:Presentation):
    """Deletes the template slides from the presentation. (the first ~7 slides)"""
    for _ in range(len(TEMPLATE)):
        presentation.slides._sldIdLst.remove(presentation.slides._sldIdLst[0])

def insert_logo_slide(presentation:Presentation):
    """Inserts the logo slide."""
    duplicate_slide(presentation, TEMPLATE.index("logo_slide"))

def insert_video_slide(presentation:Presentation, video_path:str, caption:str=""):
    """Inserts a slide with a video and a caption."""
    slide = duplicate_slide(presentation, TEMPLATE.index("media_with_caption_slide"))
    caption = caption.strip()

    # getting size and aspect ratio of the video
    clip = VideoFileClip(video_path)
    aspect_ratio = clip.size[0] / clip.size[1]
    height = presentation.slide_height.inches
    width = presentation.slide_height.inches * aspect_ratio
    left = (presentation.slide_width.inches - width) / 2

    # changing the caption
    if caption != "":
        for shape in slide.shapes:
            if shape.name == "caption":
                shape.text_frame.clear()
                caption_paragraph = shape.text_frame.paragraphs[0]
                caption_paragraph.alignment = PP_ALIGN.CENTER
                insert_text(caption_paragraph, caption, size=0, colored=False)
    else:
        # removing the caption shape if there's no caption
        for shape in slide.shapes:
            if shape.name == "caption":
                slide.shapes._spTree.remove(shape._element)

    # addding the video
    video = slide.shapes.add_movie(video_path, Inches(left), 0, Inches(width), Inches(height), poster_frame_image="gyp-thumbnail.png")
    slide.shapes._spTree.remove(video._element)
    slide.shapes._spTree.insert(2, video._element)
    autoplay_media(video)

def insert_image_slide(presentation:Presentation, image_path:str, caption:str=""):
    """Inserts a slide with an image and a caption."""
    slide = duplicate_slide(presentation, TEMPLATE.index("media_with_caption_slide"))
    caption = caption.strip()

    # getting size and aspect ratio of the image
    img = Image.open(image_path)
    width, height = img.size
    aspect_ratio = width / height
    height = presentation.slide_height.inches
    width = presentation.slide_height.inches * aspect_ratio
    left = (presentation.slide_width.inches - width) / 2

    # changing the caption
    if caption != "":
        for shape in slide.shapes:
            if shape.name == "caption":
                shape.text_frame.clear()
                caption_paragraph = shape.text_frame.paragraphs[0]
                caption_paragraph.alignment = PP_ALIGN.CENTER
                insert_text(caption_paragraph, caption, size=0, colored=False)
    else:
        # removing the caption shape if there's no caption
        for shape in slide.shapes:
            if shape.name == "caption":
                slide.shapes._spTree.remove(shape._element)

    # adding the image
    image = slide.shapes.add_picture(image_path, Inches(left), 0, Inches(width), Inches(height))
    slide.shapes._spTree.remove(image._element)
    slide.shapes._spTree.insert(2, image._element)

def insert_simple_title_slide(presentation:Presentation, title:str):
    """Inserts a slide with a simple title."""
    slide = duplicate_slide(presentation, TEMPLATE.index("simple_title_slide"))

    # changing the contents of the slide
    for shape in slide.shapes:
        if shape.name == "title":
            shape.text_frame.clear()
            title_paragraph = shape.text_frame.paragraphs[0]
            title_paragraph.alignment = PP_ALIGN.CENTER
            insert_text(title_paragraph, title, size=2, colored=False)

def insert_song_title_slide(presentation:Presentation, number:int, title:str):
    """Inserts a slide with the title of the song."""
    number = str(f"{number:03}")
    slide = duplicate_slide(presentation, TEMPLATE.index("song_title_slide"))

    # changing the contents of the slide
    for shape in slide.shapes:
        if shape.name == "song_title":
            shape.text_frame.clear()
            title_paragraph = shape.text_frame.paragraphs[0]
            title_paragraph.alignment = PP_ALIGN.CENTER
            insert_text(title_paragraph, f"#{number}", size=0, colored=False)
            insert_text(title_paragraph, f"\n{title}", size=2, colored=False)

def insert_chorus_slide(presentation:Presentation, verse_name:str, text:str, last_slide:bool=False):
    """Inserts a slide with the chorus of the song."""
    slide = duplicate_slide(presentation, TEMPLATE.index("chorus_slide"))

    # changing the contents of the slide
    for shape in slide.shapes:
        if shape.name == "chorus":
            shape.text_frame.clear()
            chorus_paragraph = shape.text_frame.paragraphs[0]
            chorus_paragraph.alignment = PP_ALIGN.CENTER
            insert_text(chorus_paragraph, verse_name, size=0, colored=True)
            insert_text(chorus_paragraph, f"\n{text}", size=1, colored=False)

    if last_slide:
        # get size of the slide in inches
        posx = Inches(presentation.slide_width.inches) - Inches(0.8)
        posy = Inches(presentation.slide_height.inches) - Inches(0.8)

        # add the square
        square = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, posx, posy, Inches(0.4), Inches(0.4))
        square.fill.solid()
        square.fill.fore_color.rgb = RGBColor(255, 255, 255)
        # remove outline
        square.line.fill.background()

def insert_scripture_slide(presentation:Presentation, reference:str, text:tuple, verse_separator:str=" "):
    """Inserts a slide with the verses.
    Text need to be provided like this:
    (
        ("1", "This is the first verse."),
        ("2", "This is the second verse."),
        ("3", "This is the third verse.")
    )
    """
    slide = duplicate_slide(presentation, TEMPLATE.index("scripture_slide"))

    # changing the contents of the slide
    for shape in slide.shapes:
        if shape.name == "scripture":
            shape.text_frame.clear()
            verse_paragraph = shape.text_frame.paragraphs[0]
            verse_paragraph.alignment = PP_ALIGN.CENTER
            for verse in text:
                if verse[0] != "":
                    if verse[0] != text[0][0]:
                        insert_text(verse_paragraph, verse_separator, size=1, colored=False)
                    insert_text(verse_paragraph, verse[0]+' ', size=1, colored=True, superscript=True)
                insert_text(verse_paragraph, verse[1], size=1, colored=False)
        elif shape.name == "reference":
            shape.text_frame.clear()
            reference_paragraph = shape.text_frame.paragraphs[0]
            reference_paragraph.alignment = PP_ALIGN.CENTER
            insert_text(reference_paragraph, reference, size=0, colored=True)

def insert_end_slide(presentation:Presentation):
    """Inserts the end slide."""
    duplicate_slide(presentation, TEMPLATE.index("end_slide"))
