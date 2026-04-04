"""Utilities to merge .pptx files into an existing presentation."""


from pptx import Presentation

from slide_duplication_utility import (
    _object_rels,
    _exp_add_slide,
    copy_shapes as copy_shapes_from_util,
    remove_shape,
)


def merge_pptx(target: Presentation, source_path: str) -> None:
    """
    Merge all slides from a source .pptx file into a target presentation.

    :param target: The target Presentation to merge into
    :param source_path: Path to the source .pptx file to merge
    """
    source = Presentation(source_path)

    for src_slide in source.slides:
        _merge_slide(target, src_slide)


def _merge_slide(target: Presentation, source_slide):
    """
    Copy a single slide from source to target presentation.
    Uses the same approach as duplicate_slide for consistency.
    """
    # Add a new slide to target using source slide's layout
    dest = _exp_add_slide(target, source_slide.slide_layout)

    # Remove all shapes from the default layout
    for shape in dest.shapes:
        remove_shape(shape)

    # Copy all existing shapes using the existing utility
    copy_shapes_from_util(source_slide.shapes, dest)

    # Copy existing references of known type (hyperlinks, etc.)
    known_refs = [
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
    ]
    for rel in _object_rels(source_slide.part):
        if rel.reltype in known_refs:
            if rel.is_external:
                dest.part.rels.get_or_add_ext_rel(rel.reltype, rel._target)
            else:
                dest.part.rels.get_or_add(rel.reltype, rel._target)

    # Copy notes if present
    if source_slide.has_notes_slide:
        txt = source_slide.notes_slide.notes_text_frame.text
        dest.notes_slide.notes_text_frame.text = txt