"""
Me trying to learn how to generate powerpoint slides using python
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from utilities import *


# Create a presentation object
prs = Presentation("template.pptx")

# # select the first slide
# first_slide = prs.slides[0]
# first_slide_title = first_slide.shapes.title
# # change font of the title
# first_slide_title_text_frame = first_slide_title.text_frame
# run = first_slide_title_text_frame.paragraphs[0].add_run()
# run.text = "LAKJKLJLKSJDKLAJ?"
# font = run.font
# font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)

# # Add a slide to the presentation
# my_slide_layout = prs.slide_layouts[TITLE_LAYOUT]
# slide = prs.slides.add_slide(my_slide_layout)

# # Add text to the slide
# title = slide.shapes.title
# title.text = "asdasdasd, World!"

# # Add a subtitle to the slide
# subtitle = slide.placeholders[1]
# subtitle.text = "This is a subtitle"

insert_video_slide(prs, "gyp.mp4", "Never Would I")
insert_song_title_slide(prs, 1, "Never Would I")
insert_image_slide(prs, "gyp-thumbnail.png", "")
insert_logo_slide(prs)
insert_simple_title_slide(prs, "This is a simple title slide")

delete_template_slides(prs)


# Save the presentation
prs.save("output.pptx")